"""
=============================================================================
Generate Project Presentation (.pptx)
=============================================================================
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

OUTPUT_DIR = os.path.dirname(os.path.abspath(__file__))
OUTPUT_FILE = os.path.join(OUTPUT_DIR, 'Project_Presentation.pptx')

DARK_BLUE = RGBColor(0x2E, 0x86, 0xAB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_BG = RGBColor(0x1E, 0x1E, 0x1E)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)


def add_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color=DARK_BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height, size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return txBox


def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG)
    add_shape(slide, Inches(1), Inches(2), Inches(11.333), Inches(3.5), DARK_BLUE)
    add_text(slide, "E-Commerce Sales Analytics", Inches(2), Inches(2.3), Inches(9.333), Inches(1.5), size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "Business Intelligence Dashboard", Inches(2), Inches(3.5), Inches(9.333), Inches(1), size=28, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "A Complete Business Analyst Portfolio Project", Inches(2), Inches(4.5), Inches(9.333), Inches(0.8), size=18, color=RGBColor(0x44, 0xBB, 0xA4), align=PP_ALIGN.CENTER)

    # --- Slide 2: Project Overview ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG)
    add_text(slide, "Project Overview", Inches(0.5), Inches(0.3), Inches(12), Inches(1), size=36, bold=True, color=DARK_BLUE)
    items = [
        "Analyzed 15,000+ e-commerce transactions to drive business decisions",
        "Built normalized PostgreSQL database with 11 tables and 50+ SQL queries",
        "Developed Python ETL pipeline for data cleaning, validation & analysis",
        "Created 8-page interactive Power BI dashboard with 49+ visuals",
        "Generated automated Excel/CSV reporting with professional formatting",
        "Produced comprehensive documentation (BRD, FS, Technical Docs, Data Dictionary)"
    ]
    for i, item in enumerate(items):
        add_text(slide, f"  {i+1}. {item}", Inches(0.5), Inches(1.5 + i * 0.8), Inches(12), Inches(0.7), size=16, color=WHITE)

    # --- Slide 3: Tech Stack ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG)
    add_text(slide, "Technology Stack", Inches(0.5), Inches(0.3), Inches(12), Inches(1), size=36, bold=True, color=DARK_BLUE)
    techs = [("Data Processing", "Python, Pandas, NumPy"), ("Database", "PostgreSQL (11 Tables, Views, SPs)"), ("Visualization", "Power BI, Matplotlib, Seaborn"), ("Reporting", "Excel (xlsxwriter), CSV"), ("Tools", "VS Code, Git, GitHub")]
    for i, (cat, tools) in enumerate(techs):
        add_shape(slide, Inches(0.5 + i * 2.5), Inches(2), Inches(2.3), Inches(4), RGBColor(0x25, 0x25, 0x26))
        add_text(slide, cat, Inches(0.7 + i * 2.5), Inches(2.3), Inches(2), Inches(0.6), size=20, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
        add_text(slide, tools, Inches(0.7 + i * 2.5), Inches(3.2), Inches(2), Inches(2), size=14, color=WHITE, align=PP_ALIGN.CENTER)

    # --- Slide 4: KPI Dashboard ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG)
    add_text(slide, "Key Performance Indicators", Inches(0.5), Inches(0.3), Inches(12), Inches(1), size=36, bold=True, color=DARK_BLUE)
    kpis = [("Total Sales", "$7.5M"), ("Total Profit", "$2.7M"), ("Profit Margin", "35.6%"), ("Total Orders", "15,000"), ("Avg Order Value", "$502"), ("Total Customers", "2,000"), ("Repeat Rate", "99.7%"), ("Top Category", "Technology"), ("Return Rate", "9.9%")]
    for i, (kpi, val) in enumerate(kpis):
        col = i % 3
        row = i // 3
        x = 0.5 + col * 4.2
        y = 2 + row * 1.8
        add_shape(slide, Inches(x), Inches(y), Inches(3.8), Inches(1.5), RGBColor(0x25, 0x25, 0x26))
        add_text(slide, kpi, Inches(x + 0.2), Inches(y + 0.2), Inches(3.4), Inches(0.5), size=18, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
        add_text(slide, val, Inches(x + 0.2), Inches(y + 0.7), Inches(3.4), Inches(0.5), size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # --- Slide 5: SQL Features ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG)
    add_text(slide, "SQL Analytics Features", Inches(0.5), Inches(0.3), Inches(12), Inches(1), size=36, bold=True, color=DARK_BLUE)
    sql_features = ["50+ Analytical SQL Queries", "Window Functions for Running Totals", "CTEs for Complex Analysis", "Views for Reusable Reports", "Stored Procedures for Automation", "Performance-Optimized Indexes", "Revenue, Profit, Customer, Product Analysis", "Year-over-Year & Month-over-Month Growth"]
    for i, feature in enumerate(sql_features):
        col = i % 2
        row = i // 2
        x = 0.5 + col * 6.3
        y = 1.5 + row * 1.5
        add_shape(slide, Inches(x), Inches(y), Inches(5.8), Inches(1.2), RGBColor(0x25, 0x25, 0x26))
        add_text(slide, f"  {feature}", Inches(x + 0.2), Inches(y + 0.2), Inches(5.4), Inches(0.8), size=18, color=WHITE)

    # --- Slide 6: Dashboard Pages ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG)
    add_text(slide, "Power BI Dashboard - 8 Pages", Inches(0.5), Inches(0.3), Inches(12), Inches(1), size=36, bold=True, color=DARK_BLUE)
    pages = [("Executive Summary", "KPI Cards, Trend Lines, US Map"), ("Sales Dashboard", "Bar Charts, Pie Charts, Area Chart"), ("Profit Dashboard", "Scatter Plot, Margin Analysis"), ("Customer Dashboard", "Segmentation, LTV Analysis"), ("Regional Dashboard", "Geographic Map, State Analysis"), ("Product Dashboard", "Treemap, Category Performance"), ("Forecast Dashboard", "Predictive Analytics, Trends"), ("KPI Dashboard", "Gauges, Target vs Actual")]
    for i, (page, desc) in enumerate(pages):
        col = i % 4
        row = i // 4
        x = 0.3 + col * 3.25
        y = 1.5 + row * 3
        add_shape(slide, Inches(x), Inches(y), Inches(3), Inches(2.5), RGBColor(0x25, 0x25, 0x26))
        add_text(slide, f"Page {i+1}", Inches(x + 0.2), Inches(y + 0.2), Inches(2.6), Inches(0.4), size=12, color=RGBColor(0x44, 0xBB, 0xA4), align=PP_ALIGN.CENTER)
        add_text(slide, page, Inches(x + 0.2), Inches(y + 0.6), Inches(2.6), Inches(0.5), size=20, bold=True, color=DARK_BLUE, align=PP_ALIGN.CENTER)
        add_text(slide, desc, Inches(x + 0.2), Inches(y + 1.3), Inches(2.6), Inches(1), size=14, color=WHITE, align=PP_ALIGN.CENTER)

    # --- Slide 7: Business Impact ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG)
    add_text(slide, "Business Impact & Insights", Inches(0.5), Inches(0.3), Inches(12), Inches(1), size=36, bold=True, color=DARK_BLUE)
    insights = [
        "Identified Technology as top revenue category ($3.9M)",
        "South region leads in sales; West region in profit margin",
        "Consumer segment highest volume; Corporate highest AOV",
        "Discounted orders show 15% higher volume but 8% lower margin",
        "Q4 holiday season shows 40% sales spike vs Q1 average",
        "Same Day shipping has highest profit margin despite cost",
        "Proposed loyalty program to improve customer retention",
        "Recommended discount optimization to protect margins"
    ]
    for i, insight in enumerate(insights):
        add_text(slide, f"  \u2022 {insight}", Inches(0.5), Inches(1.5 + i * 0.7), Inches(12), Inches(0.6), size=15, color=WHITE)

    # --- Slide 8: Thank You ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide, DARK_BG)
    add_shape(slide, Inches(3), Inches(2.5), Inches(7.333), Inches(3), DARK_BLUE)
    add_text(slide, "Thank You", Inches(3.5), Inches(2.8), Inches(6.333), Inches(1.5), size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(slide, "GitHub Ready  |  ATS Friendly  |  Resume Worthy", Inches(3.5), Inches(4.2), Inches(6.333), Inches(0.8), size=20, color=RGBColor(0x44, 0xBB, 0xA4), align=PP_ALIGN.CENTER)

    prs.save(OUTPUT_FILE)
    print(f"Presentation saved: {OUTPUT_FILE}")
    return OUTPUT_FILE


if __name__ == '__main__':
    create_presentation()
